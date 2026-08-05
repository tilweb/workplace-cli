"""Document-builder skill — create Word/Excel/PowerPoint/PDF files.

Guidance skill (no on-disk resources): the agent writes a short Python
script using the bundled libraries and runs it via the ``bash`` tool. The
libraries (python-docx, openpyxl, python-pptx, reportlab) ship with
Workplace CLI, so they are always importable in the tool's own environment.

User-invocable via ``/document-builder`` in the TUI.
"""

from __future__ import annotations

from vibe.core.skills.models import SkillInfo

_DESCRIPTION = (
    "Create office documents — Word (.docx), Excel (.xlsx), PowerPoint (.pptx) "
    "and PDF — from user-provided content or data. Use this skill whenever the "
    "user asks to generate, export, or 'make me' a report, letter, spreadsheet, "
    "table, slide deck, or PDF. Produces clean, neutral, standard-looking "
    "documents by writing and running a short Python script."
)

_PROMPT = r"""\
This skill creates real office documents (.docx, .xlsx, .pptx, .pdf) by
writing a short Python script and running it with the `bash` tool.

**CRITICAL — which Python to use.** The libraries below ship *with Workplace
CLI*, installed in its own virtual environment. The shell's `python`/`python3`
is the user's *system* interpreter (e.g. Homebrew) and does NOT have them —
running `python3 script.py` fails with `ModuleNotFoundError`, and
`pip install` is both wrong and usually blocked. Always run your script with
the interpreter Workplace CLI exposes via the `$WORKPLACE_PYTHON` environment
variable, which points at the venv that already has every library:

```bash
"$WORKPLACE_PYTHON" build_doc.py
```

Never `pip install` these and never fall back to bare `python`/`python3`:

- **Word (.docx)** → `python-docx`   (`import docx`)
- **Excel (.xlsx)** → `openpyxl`      (`import openpyxl`)
- **PowerPoint (.pptx)** → `python-pptx` (`import pptx`)
- **PDF** → `reportlab`               (`from reportlab.platypus import ...`)

## Workflow

1. **Pick the format.** If the user named one (Word, Excel, PDF, slides), use
   it. Otherwise choose the natural fit: prose/letters/reports → Word or PDF;
   tabular data → Excel; presentations → PowerPoint. When unsure, ask.
2. **Decide the output path.** Default to the current working directory with a
   sensible file name (e.g. `angebot.docx`). Respect any path the user gives.
3. **Write a script**, then run it with `"$WORKPLACE_PYTHON" script.py` (see
   above — never bare `python`/`python3`). Prefer writing the script to the
   scratchpad dir and running it, over a giant inline `-c`. Keep the content
   in the script; do not hardcode secrets.
4. **Verify**: confirm the file exists and report its absolute path and size.
   For a quick sanity check you can `read_file` the result (docx/pdf render as
   text/pages). Never claim success without the file on disk.

## Design: neutral and standard

The user wants clean, professional, unobtrusive documents — NOT flashy design.
Follow these defaults unless the user asks otherwise:

- Keep the library's **default template, fonts, and colors**. Standard body
  font (Word: Calibri/Aptos default; PDF: Helvetica), ~11pt body, dark-grey or
  black text on white.
- Use **built-in heading styles** for structure, not manual bold+size hacks.
- Tables: thin single-line borders, a bold header row, left-aligned text,
  right-aligned numbers. No heavy fills or accent colors.
- Consistent margins, generous white space, no decorative graphics.
- Only add color/branding when the user explicitly asks.

## Minimal patterns

### Word (.docx)
```python
from docx import Document
from docx.shared import Pt

doc = Document()  # default template = neutral
doc.add_heading("Projektbericht", level=0)
doc.add_heading("Zusammenfassung", level=1)
doc.add_paragraph("Fließtext im Standard-Stil.")

rows = [("Posten", "Menge", "Preis"), ("Lizenz", "3", "1.200 €")]
table = doc.add_table(rows=len(rows), cols=3)
table.style = "Table Grid"  # neutral thin borders
for r, cells in enumerate(rows):
    for c, val in enumerate(cells):
        table.cell(r, c).text = str(val)
for cell in table.rows[0].cells:      # bold header row
    cell.paragraphs[0].runs[0].font.bold = True

doc.save("bericht.docx")
```

### Excel (.xlsx)
```python
from openpyxl import Workbook
from openpyxl.styles import Font

wb = Workbook()
ws = wb.active
ws.title = "Daten"
ws.append(["Datum", "Kunde", "Betrag"])
for cell in ws[1]:
    cell.font = Font(bold=True)       # bold header
ws.append(["2026-08-05", "Adacor", 1200])
ws.freeze_panes = "A2"                # keep header visible
for col in ws.columns:                # simple auto width
    width = max(len(str(c.value)) for c in col if c.value is not None) + 2
    ws.column_dimensions[col[0].column_letter].width = width
wb.save("daten.xlsx")
```

### PowerPoint (.pptx)
```python
from pptx import Presentation

prs = Presentation()  # default 4:3 neutral template
title = prs.slides.add_slide(prs.slide_layouts[0])   # title layout
title.shapes.title.text = "Quartalsupdate"
title.placeholders[1].text = "Adacor · Q3 2026"

body = prs.slides.add_slide(prs.slide_layouts[1])    # title + bullets
body.shapes.title.text = "Highlights"
tf = body.placeholders[1].text_frame
tf.text = "Erster Punkt"
tf.add_paragraph().text = "Zweiter Punkt"

prs.save("deck.pptx")
```

### PDF
```python
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

styles = getSampleStyleSheet()
doc = SimpleDocTemplate("bericht.pdf", pagesize=A4)
story = [
    Paragraph("Projektbericht", styles["Title"]),
    Spacer(1, 12),
    Paragraph("Standard-Fließtext in Helvetica.", styles["BodyText"]),
    Spacer(1, 12),
]
data = [["Posten", "Preis"], ["Lizenz", "1.200 €"]]
table = Table(data)
table.setStyle(TableStyle([
    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),   # header
    ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
]))
story.append(table)
doc.build(story)
```

## Notes

- For non-ASCII text (German umlauts ä/ö/ü/ß, €, em-dash —), all four
  libraries handle it out of the box — just use normal Python strings. In
  particular reportlab's default fonts (Helvetica/Times) already cover
  Latin-1, so do NOT register an external TTF (e.g. DejaVu) for German text —
  those font paths often don't exist on macOS and the registration fails.
  Only register a font for scripts outside Latin-1 (e.g. CJK).
- Large tables/data: read the source (CSV, prior tool output, user text) first,
  then build the document from it. For Excel, write numbers as numbers (not
  strings) so they stay sortable and summable.
- If the user wants an existing file edited, open it (`Document(path)`,
  `openpyxl.load_workbook(path)`, `Presentation(path)`) instead of starting
  fresh.
"""

SKILL = SkillInfo(
    name="document-builder",
    description=_DESCRIPTION,
    user_invocable=True,
    prompt=_PROMPT,
)
